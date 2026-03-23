"""Generate a PowerPoint presentation (.pptx) from an OrderOfWorship."""

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

from .models import OrderOfWorship, HymnRef
from .scripture import fetch_scripture
from .themes import get_theme

RESOURCES = Path(__file__).parent.parent.parent / "resources"
HYMNAL_DIR = Path(__file__).parent.parent.parent / "hymnal-json"
OUTPUT_DIR = Path(__file__).parent.parent.parent / "output"

# Slide dimensions: 13.333" x 7.500" (widescreen 16:9)
SLIDE_WIDTH = Emu(12_192_000)
SLIDE_HEIGHT = Emu(6_858_000)

# Backgrounds
HYMN_BG = RESOURCES / "slides" / "Background.png"
CREED_BG = RESOURCES / "slides" / "CreedBackground.png"

# Static slide images
ANNOUNCEMENTS_SLIDE = RESOURCES / "slides" / "AnnouncementsSlide.jpg"
CONCERNS_SLIDE = RESOURCES / "slides" / "ConcernsSlide.jpg"
OFFERING_SLIDE = RESOURCES / "slides" / "OfferingSlide.jpg"
PRAYER_SLIDE = RESOURCES / "slides" / "PrayerSlide.jpg"

# Active theme — loaded at generation time, module-level default for helpers
_theme = get_theme()

# Theme accessors (set before each generation via _load_theme)
THEME_FONT = _theme["font"]
THEME_TITLE_COLOR = _theme["title_color"]
THEME_TEXT_COLOR = _theme["text_color"]
THEME_BADGE_BG = _theme["badge_bg"]
THEME_BADGE_FG = _theme["badge_fg"]
THEME_SPEAKER_LABEL_COLOR = _theme["speaker_label_color"]
THEME_LITURGY_UPPERCASE = _theme["liturgy_uppercase"]
THEME_SHADOW_ENABLED = _theme["shadow_enabled"]


def _load_theme(theme_name: str = None):
    """Load a theme and update module-level variables."""
    global THEME_FONT, THEME_TITLE_COLOR, THEME_TEXT_COLOR, THEME_BADGE_BG, THEME_BADGE_FG
    global THEME_SPEAKER_LABEL_COLOR, THEME_LITURGY_UPPERCASE, THEME_SHADOW_ENABLED
    t = get_theme(theme_name)
    THEME_FONT = t["font"]
    THEME_TITLE_COLOR = t["title_color"]
    THEME_TEXT_COLOR = t["text_color"]
    THEME_BADGE_BG = t["badge_bg"]
    THEME_BADGE_FG = t["badge_fg"]
    THEME_SPEAKER_LABEL_COLOR = t["speaker_label_color"]
    THEME_LITURGY_UPPERCASE = t["liturgy_uppercase"]
    THEME_SHADOW_ENABLED = t["shadow_enabled"]

# Non-theme colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Layout constants (from analysis of existing presentation)
TEXT_LEFT = Emu(1_524_000)           # 1.667 inches
TEXT_WIDTH = Emu(9_144_000)          # 10.000 inches
TITLE_TOP = Emu(304_800)            # 0.333 inches
BG_IMG_LEFT = Emu(2_667_000)        # 2.917 inches
BG_IMG_SIZE = Emu(6_858_000)        # 7.500 inches (square)

# Text margins
MARGIN_LR = Emu(91_440)
MARGIN_TB = Emu(45_720)


def _add_shadow(shape):
    """Add outer shadow effect to a shape for text readability."""
    if not THEME_SHADOW_ENABLED:
        return
    from lxml import etree
    spPr = shape._element.spPr
    if spPr is None:
        from pptx.oxml.ns import qn
        spPr = etree.SubElement(shape._element, qn('a:spPr'))

    # Build effectLst with outer shadow
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    effectLst = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')
    outerShdw = etree.SubElement(effectLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw')
    outerShdw.set('dist', '35921')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'ctr')
    outerShdw.set('rotWithShape', '0')
    schemeClr = etree.SubElement(outerShdw, '{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
    schemeClr.set('val', 'bg1')


def _set_paragraph_spacing(paragraph, space_before_pct=50):
    """Set paragraph spacing before as percentage of font size."""
    from lxml import etree
    pPr = paragraph._p.get_or_add_pPr()
    spcBef = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcBef')
    spcPct = etree.SubElement(spcBef, '{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
    spcPct.set('val', str(space_before_pct * 1000))


def _add_full_image_slide(prs, image_path):
    """Add a slide with a single full-bleed image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path), Emu(0), Emu(0),
            SLIDE_WIDTH, SLIDE_HEIGHT
        )
    return slide


def _add_hymn_background(slide):
    """Add the hymn background image (7.5x7.5 square, right-aligned)."""
    if HYMN_BG.exists():
        slide.shapes.add_picture(
            str(HYMN_BG), BG_IMG_LEFT, Emu(0),
            BG_IMG_SIZE, BG_IMG_SIZE
        )


def _add_creed_background(slide):
    """Add the creed/liturgy background image (full-bleed)."""
    if CREED_BG.exists():
        slide.shapes.add_picture(
            str(CREED_BG), Emu(0), Emu(0),
            SLIDE_WIDTH, SLIDE_HEIGHT
        )


def _load_hymn_data(ref: HymnRef) -> dict | None:
    """Load hymn JSON data from the hymnal-json directory."""
    source_dir = HYMNAL_DIR / ref.source
    # Try to find the file by number prefix
    for f in source_dir.glob("*.json"):
        data = json.load(open(f, encoding='utf-8'))
        if data.get('number') == ref.number:
            return data
    return None


def _is_attribution_line(line: str) -> bool:
    """Check if a line is a copyright attribution."""
    line_stripped = line.strip()
    if '©' in line_stripped or '\u00a9' in line_stripped:
        return True
    if line_stripped.startswith('FROM THE'):
        return True
    if line_stripped.startswith('FORMER ') and line_stripped.isupper():
        return True
    return False


def _is_refrain_label(line: str) -> bool:
    """Check if a line is just 'Refrain'."""
    return line.strip().lower() == 'refrain'


def _parse_hymn_slides(hymn_data: dict) -> list[dict]:
    """
    Parse hymn JSON into structured slide data.

    The JSON data is already clean — the extraction script uses
    position-based shape identification to separate title/number/attribution
    from lyrics at the shape level. So we don't need to do title-matching
    here. We only need to identify copyright lines and refrain labels.

    Returns list of dicts with keys:
      - type: 'first' | 'continuation' | 'refrain'
      - title: str (for first slides)
      - number: str (for first slides)
      - attribution: str (for first slides)
      - lyrics: list[str]
      - refrain: bool
      - verse_label: str (e.g. "(Verse 1)")
    """
    slides = hymn_data['slides']
    title = hymn_data['title']
    number = hymn_data['number']

    # Shorten verbose Lord's Prayer titles at runtime
    if number in ('894', '895', '896'):
        title = "The Lord's Prayer"
    result = []

    for si, slide in enumerate(slides):
        lines = slide['lines']

        # Separate attribution/copyright from lyrics
        lyrics = []
        attribution = ""
        verse_label = ""
        is_refrain = False
        is_first = (si == 0)

        for line in lines:
            if _is_attribution_line(line):
                attribution = line.strip()
                continue
            if re.match(r'^\((?:verse\s+)?\d+\)$', line.strip(), re.IGNORECASE):
                verse_label = line.strip()
                continue
            if line.strip().lower() == '(refrain)':
                verse_label = "(Refrain)"
                continue
            if _is_refrain_label(line):
                is_refrain = True
                continue
            lyrics.append(line)

        slide_info = {
            'type': 'first' if is_first else ('refrain' if is_refrain else 'continuation'),
            'title': title if is_first else '',
            'number': number if is_first else '',
            'source': hymn_data.get('source', ''),
            'attribution': attribution,
            'lyrics': lyrics,
            'refrain': is_refrain,
            'verse_label': verse_label,
        }
        result.append(slide_info)

    return result


def _source_label(source: str, number: str) -> str:
    """Build a source tag like 'UMH 64' or 'TFWS 2001'."""
    # Map folder names to hymnal display names
    source_map = {
        'umh': 'UMH',
        'umh-services': 'UMH',
        'umh-general': 'UMH',
        'tfws': 'TFWS',
    }
    prefix = source_map.get(source.lower(), source.upper()) if source else ""
    # Strip leading zeros for cleaner display
    num = number.lstrip('0') or number
    return f"{prefix} {num}".strip()


def _add_source_badge(slide, label: str):
    """Add a modern semi-transparent pill badge in the bottom-right corner."""
    from pptx.oxml.ns import qn
    from lxml import etree
    from pptx.enum.shapes import MSO_SHAPE

    badge_width = Emu(2_100_000)   # ~2.3 inches
    badge_height = Emu(500_000)    # ~0.55 inches
    badge_left = SLIDE_WIDTH - badge_width - Emu(274_320)   # 0.3" from right edge
    badge_top = SLIDE_HEIGHT - badge_height - Emu(274_320)  # 0.3" from bottom

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        badge_left, badge_top, badge_width, badge_height
    )

    # Rounded corners — adjust the radius
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 30000')  # ~30% corner radius for pill shape

    # Semi-transparent dark fill
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = THEME_BADGE_BG

    # Set 30% transparency via XML
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        solidFill = shape._element.spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '70000')  # 70% opacity (30% transparent)

    # No border
    shape.line.fill.background()

    # Dark drop shadow on the badge
    if THEME_SHADOW_ENABLED:
        from lxml import etree as _et
        from pptx.oxml.ns import qn as _qn
        _spPr = shape._element.spPr
        _effectLst = _et.SubElement(_spPr, _qn('a:effectLst'))
        _outerShdw = _et.SubElement(_effectLst, _qn('a:outerShdw'))
        _outerShdw.set('dist', '50800')     # ~4pt offset
        _outerShdw.set('dir', '5400000')    # straight down
        _outerShdw.set('blurRad', '76200')  # ~6pt blur
        _outerShdw.set('algn', 'ctr')
        _outerShdw.set('rotWithShape', '0')
        _srgb = _et.SubElement(_outerShdw, _qn('a:srgbClr'))
        _srgb.set('val', '000000')
        _alpha = _et.SubElement(_srgb, _qn('a:alpha'))
        _alpha.set('val', '50000')  # 50% opacity black

    # Text
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = THEME_FONT
    run.font.size = Pt(24)
    run.font.color.rgb = THEME_BADGE_FG
    run.font.bold = True

    # Letter spacing for a clean feel
    rPr = run._r.get_or_add_rPr()
    spc = etree.SubElement(rPr, qn('a:spc'))
    spc.set('val', '200')  # 2pt letter spacing


def _add_title_pill(slide, title_text: str, top=None, font_size=48):
    """Add a title in a themed pill shape — same style as source badge but for titles."""
    from pptx.oxml.ns import qn
    from lxml import etree
    from pptx.enum.shapes import MSO_SHAPE

    if top is None:
        top = TITLE_TOP

    # Pill stretches most of the slide width, centered
    pill_margin = Emu(800_000)  # ~0.87" each side
    pill_left = pill_margin
    pill_width = SLIDE_WIDTH - (pill_margin * 2)

    # Balanced line splitting for long titles (split at word boundary nearest midpoint)
    chars_per_line = 28 if font_size >= 48 else 35 if font_size >= 40 else 45
    if len(title_text) > chars_per_line:
        words = title_text.split()
        if len(words) > 2:
            best_split = None
            best_diff = float('inf')
            running = ''
            for i, word in enumerate(words[:-1]):
                running += ('' if i == 0 else ' ') + word
                rest = ' '.join(words[i + 1:])
                diff = abs(len(running) - len(rest))
                if diff < best_diff:
                    best_diff = diff
                    best_split = (running, rest)
            if best_split:
                title_text = best_split[0] + '\n' + best_split[1]

    # Calculate line count for manual height (AutoShape ignores auto_size)
    if '\n' in title_text:
        num_lines = sum(1 + max(0, len(part) // chars_per_line) for part in title_text.split('\n'))
    else:
        num_lines = 1 + max(0, len(title_text) // chars_per_line)

    line_height = Emu(730_000)   # ~0.8" per line at 48pt
    padding = Emu(250_000)       # top + bottom padding
    pill_height = (line_height * num_lines) + padding

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        pill_left, top, pill_width, pill_height
    )

    # Rounded corners
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 25000')  # slightly less round than badge

    # Same semi-transparent fill as badge
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = THEME_BADGE_BG
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        solidFill = shape._element.spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '70000')

    # No border
    shape.line.fill.background()

    # Dark drop shadow
    if THEME_SHADOW_ENABLED:
        _effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        _outerShdw = etree.SubElement(_effectLst, qn('a:outerShdw'))
        _outerShdw.set('dist', '50800')
        _outerShdw.set('dir', '5400000')
        _outerShdw.set('blurRad', '76200')
        _outerShdw.set('algn', 'ctr')
        _outerShdw.set('rotWithShape', '0')
        _srgb = etree.SubElement(_outerShdw, qn('a:srgbClr'))
        _srgb.set('val', '000000')
        _alpha = etree.SubElement(_srgb, qn('a:alpha'))
        _alpha.set('val', '50000')

    # Text (no auto_size — AutoShapes ignore SHAPE_TO_FIT_TEXT)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(182_880)   # ~0.2"
    tf.margin_right = Emu(182_880)
    tf.margin_top = Emu(91_440)     # ~0.1"
    tf.margin_bottom = Emu(91_440)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = title_text
    p.font.name = THEME_FONT
    p.font.size = Pt(font_size)
    p.font.color.rgb = THEME_BADGE_FG
    p.font.bold = True

    return shape, num_lines


def _create_hymn_first_slide(prs, slide_info: dict):
    """Create a hymn's first slide with title pill, attribution, lyrics, background, and source badge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # 1. Title pill
    title_text = slide_info['title']
    _, title_lines = _add_title_pill(slide, title_text, top=TITLE_TOP, font_size=48)

    # Dynamic offset: base + extra per line beyond 1
    extra_lines = max(0, title_lines - 1)
    line_offset = Emu(730_000)  # ~0.8" per extra title line

    # 3. Attribution line (if present)
    attr_top = Emu(1_200_000) + (line_offset * extra_lines)
    if slide_info['attribution']:
        attr_box = slide.shapes.add_textbox(
            TEXT_LEFT, attr_top, TEXT_WIDTH, Emu(400_000)
        )
        attr_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = attr_box.text_frame.paragraphs[0]
        p.text = slide_info['attribution']
        p.font.name = THEME_FONT
        p.font.size = Pt(22)
        p.font.color.rgb = THEME_TITLE_COLOR
        _set_paragraph_spacing(p, 50)
        _add_shadow(attr_box)
        lyrics_top = Emu(1_700_000) + (line_offset * extra_lines)
    else:
        lyrics_top = Emu(1_400_000) + (line_offset * extra_lines)

    # 4. Lyrics text box
    if slide_info['lyrics']:
        lyrics_box = slide.shapes.add_textbox(
            TEXT_LEFT, lyrics_top, TEXT_WIDTH, Emu(500_000)
        )
        lyrics_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lyrics_box.text_frame.margin_left = MARGIN_LR
        lyrics_box.text_frame.margin_right = MARGIN_LR
        lyrics_box.text_frame.margin_top = MARGIN_TB
        lyrics_box.text_frame.margin_bottom = MARGIN_TB

        # Use soft returns (\x0b) to join lines within one paragraph
        p = lyrics_box.text_frame.paragraphs[0]
        lyrics_text = '\x0b'.join(slide_info['lyrics'])
        p.text = lyrics_text
        p.font.name = THEME_FONT
        p.font.color.rgb = THEME_TEXT_COLOR

        # Scale font for first slide based on available space
        num_lyric_lines = len(slide_info['lyrics'])
        available = SLIDE_HEIGHT - lyrics_top - Emu(600_000)  # reserve for badge
        line_space = available // max(num_lyric_lines, 1)
        if line_space >= Emu(685_800):      # ~0.75" per line → 50pt
            p.font.size = Pt(50)
        elif line_space >= Emu(594_360):    # ~0.65" → 44pt
            p.font.size = Pt(44)
        else:                                # tight → 38pt
            p.font.size = Pt(38)

        _set_paragraph_spacing(p, 50)
        _add_shadow(lyrics_box)

    # 5. Source badge (bottom-right pill)
    source_tag = _source_label(slide_info.get('source', ''), slide_info['number'])
    _add_source_badge(slide, source_tag)

    # 6. Background image (added last so it's on top in z-order,
    #    matching original presentation structure)
    _add_hymn_background(slide)

    return slide


def _create_hymn_continuation_slide(prs, slide_info: dict):
    """Create a hymn continuation slide (lyrics only + background)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    if slide_info['lyrics']:
        # Calculate top position based on content amount
        num_lines = len(slide_info['lyrics'])
        if num_lines <= 3:
            top = Emu(1_800_000)  # center-ish for short content
        elif num_lines <= 5:
            top = Emu(1_200_000)
        else:
            top = Emu(800_000)

        lyrics_box = slide.shapes.add_textbox(
            TEXT_LEFT, top, TEXT_WIDTH, Emu(500_000)
        )
        lyrics_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lyrics_box.text_frame.margin_left = MARGIN_LR
        lyrics_box.text_frame.margin_right = MARGIN_LR
        lyrics_box.text_frame.margin_top = MARGIN_TB
        lyrics_box.text_frame.margin_bottom = MARGIN_TB

        p = lyrics_box.text_frame.paragraphs[0]

        if slide_info['refrain']:
            # Add "Refrain" label as first paragraph, italic
            p.text = "Refrain"
            p.font.name = THEME_FONT
            p.font.size = Pt(48)
            p.font.italic = True
            p.font.color.rgb = THEME_TEXT_COLOR
            _set_paragraph_spacing(p, 50)

            # Lyrics in second paragraph
            p2 = lyrics_box.text_frame.add_paragraph()
            lyrics_text = '\x0b'.join(slide_info['lyrics'])
            p2.text = lyrics_text
            p2.font.name = THEME_FONT
            p2.font.size = Pt(50)
            p2.font.color.rgb = THEME_TEXT_COLOR
            _set_paragraph_spacing(p2, 50)
        else:
            lyrics_text = '\x0b'.join(slide_info['lyrics'])
            p.text = lyrics_text
            p.font.name = THEME_FONT
            p.font.size = Pt(50)
            p.font.color.rgb = THEME_TEXT_COLOR
            _set_paragraph_spacing(p, 50)

        _add_shadow(lyrics_box)

    _add_hymn_background(slide)
    return slide


def _create_liturgy_first_slide(prs, slide_info: dict, bg_path: Path):
    """Create a liturgy first slide (creed/Lord's Prayer) with title pill, text, and source badge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title pill
    title_text = slide_info['title'].upper() if THEME_LITURGY_UPPERCASE else slide_info['title']
    _, title_lines = _add_title_pill(slide, title_text, top=TITLE_TOP, font_size=48)
    extra_lines = max(0, title_lines - 1)
    line_offset = Emu(730_000)

    # Lyrics/text
    if slide_info['lyrics']:
        lyrics_top = Emu(1_400_000) + (line_offset * extra_lines)
        lyrics_box = slide.shapes.add_textbox(
            Emu(1_828_800), lyrics_top,       # ~2.0 inches left
            Emu(8_534_400), Emu(500_000)      # ~9.333 inches wide
        )
        lyrics_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lyrics_box.text_frame.margin_left = MARGIN_LR
        lyrics_box.text_frame.margin_right = MARGIN_LR

        # Scale font for first slide based on available space
        num_lyric_lines = len(slide_info['lyrics'])
        available = SLIDE_HEIGHT - lyrics_top - Emu(600_000)
        line_space = available // max(num_lyric_lines, 1)
        if line_space >= Emu(685_800):
            lyric_size = 48
        elif line_space >= Emu(594_360):
            lyric_size = 42
        else:
            lyric_size = 36

        p = lyrics_box.text_frame.paragraphs[0]
        _build_liturgy_text(lyrics_box.text_frame, slide_info['lyrics'], font_size=lyric_size)
        _add_shadow(lyrics_box)

    # Source badge (bottom-right pill)
    source_tag = _source_label(slide_info.get('source', ''), slide_info['number'])
    _add_source_badge(slide, source_tag)

    # Background (added last)
    if bg_path.exists():
        slide.shapes.add_picture(
            str(bg_path), Emu(0), Emu(0),
            SLIDE_WIDTH, SLIDE_HEIGHT
        )

    return slide


def _create_liturgy_continuation_slide(prs, slide_info: dict, bg_path: Path):
    """Create a liturgy continuation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if slide_info['lyrics']:
        num_lines = len(slide_info['lyrics'])
        top = Emu(800_000) if num_lines > 4 else Emu(1_400_000)

        lyrics_box = slide.shapes.add_textbox(
            Emu(1_828_800), top,
            Emu(8_534_400), Emu(500_000)
        )
        lyrics_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lyrics_box.text_frame.margin_left = MARGIN_LR
        lyrics_box.text_frame.margin_right = MARGIN_LR

        _build_liturgy_text(lyrics_box.text_frame, slide_info['lyrics'])
        _add_shadow(lyrics_box)

    if bg_path.exists():
        slide.shapes.add_picture(
            str(bg_path), Emu(0), Emu(0),
            SLIDE_WIDTH, SLIDE_HEIGHT
        )

    return slide


def _build_liturgy_text(text_frame, lines: list[str], font_size: int = 48):
    """Build liturgy text with speaker labels in red italic."""
    full_text = '\x0b'.join(lines)

    p = text_frame.paragraphs[0]

    speaker_pattern = re.compile(r'^(Pastor(?:\s+and\s+People)?|People|All|Leader|Pastor\s*:?)\s*:?\s*', re.IGNORECASE)
    has_speakers = any(speaker_pattern.match(line) for line in lines)

    if has_speakers:
        first = True
        for line in lines:
            if not first:
                from pptx.oxml.ns import qn
                from lxml import etree
                br = etree.SubElement(p._p, qn('a:br'))

            match = speaker_pattern.match(line)
            if match:
                speaker = match.group(0).strip()
                if not speaker.endswith(':'):
                    speaker += ':'
                rest = line[match.end():].strip()

                run = p.add_run()
                run.text = speaker + ' '
                run.font.name = THEME_FONT
                run.font.size = Pt(font_size)
                run.font.italic = True
                run.font.color.rgb = THEME_SPEAKER_LABEL_COLOR

                if rest:
                    run2 = p.add_run()
                    run2.text = rest
                    run2.font.name = THEME_FONT
                    run2.font.size = Pt(font_size)
                    run2.font.color.rgb = THEME_TEXT_COLOR
            else:
                run = p.add_run()
                run.text = line
                run.font.name = THEME_FONT
                run.font.size = Pt(font_size)
                run.font.color.rgb = THEME_TEXT_COLOR

            first = False
    else:
        p.text = full_text
        p.font.name = THEME_FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = THEME_TEXT_COLOR

    _set_paragraph_spacing(p, 50)


def _add_hymn_slides(prs, ref: HymnRef, bg_type: str = 'hymn'):
    """Add all slides for a hymn/creed/prayer."""
    hymn_data = _load_hymn_data(ref)
    if not hymn_data:
        return

    parsed = _parse_hymn_slides(hymn_data)
    bg_path = CREED_BG if bg_type == 'creed' else HYMN_BG

    for i, slide_info in enumerate(parsed):
        if bg_type == 'creed':
            if i == 0:
                _create_liturgy_first_slide(prs, slide_info, bg_path)
            else:
                _create_liturgy_continuation_slide(prs, slide_info, bg_path)
        else:
            if i == 0:
                _create_hymn_first_slide(prs, slide_info)
            else:
                _create_hymn_continuation_slide(prs, slide_info)


def _create_scripture_first_slide(prs, reference: str, translation_name: str, slide_data: dict):
    """Create the first scripture slide with reference title, verse text, and source badge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title pill: scripture reference (e.g., "Matthew 4:1-11")
    _, title_lines = _add_title_pill(slide, reference, top=TITLE_TOP, font_size=48)
    extra_lines = max(0, title_lines - 1)
    line_offset = Emu(730_000)

    # Verse text
    if slide_data['lines']:
        lyrics_top = Emu(1_400_000) + (line_offset * extra_lines)
        lyrics_box = slide.shapes.add_textbox(
            TEXT_LEFT, lyrics_top, TEXT_WIDTH, Emu(500_000)
        )
        lyrics_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lyrics_box.text_frame.margin_left = MARGIN_LR
        lyrics_box.text_frame.margin_right = MARGIN_LR
        lyrics_box.text_frame.margin_top = MARGIN_TB
        lyrics_box.text_frame.margin_bottom = MARGIN_TB

        p = lyrics_box.text_frame.paragraphs[0]
        lyrics_text = '\x0b'.join(slide_data['lines'])
        p.text = lyrics_text
        p.font.name = THEME_FONT
        p.font.size = Pt(44)
        p.font.color.rgb = THEME_TEXT_COLOR
        _set_paragraph_spacing(p, 50)
        _add_shadow(lyrics_box)

    # Source badge (bottom-right pill)
    _add_source_badge(slide, translation_name)

    # Background
    _add_hymn_background(slide)
    return slide


def _create_scripture_continuation_slide(prs, slide_data: dict, is_last: bool = False,
                                          translation_name: str = ""):
    """Create a scripture continuation slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    if slide_data['lines']:
        num_lines = len(slide_data['lines'])
        if num_lines <= 3:
            top = Emu(1_800_000)
        elif num_lines <= 5:
            top = Emu(1_200_000)
        else:
            top = Emu(800_000)

        lyrics_box = slide.shapes.add_textbox(
            TEXT_LEFT, top, TEXT_WIDTH, Emu(500_000)
        )
        lyrics_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        lyrics_box.text_frame.margin_left = MARGIN_LR
        lyrics_box.text_frame.margin_right = MARGIN_LR
        lyrics_box.text_frame.margin_top = MARGIN_TB
        lyrics_box.text_frame.margin_bottom = MARGIN_TB

        p = lyrics_box.text_frame.paragraphs[0]
        lyrics_text = '\x0b'.join(slide_data['lines'])
        p.text = lyrics_text
        p.font.name = THEME_FONT
        p.font.size = Pt(44)
        p.font.color.rgb = THEME_TEXT_COLOR
        _set_paragraph_spacing(p, 50)

        # Add attribution on last slide
        if is_last and translation_name:
            p2 = lyrics_box.text_frame.add_paragraph()
            p2.text = f"— {translation_name}"
            p2.font.name = THEME_FONT
            p2.font.size = Pt(18)
            p2.font.color.rgb = THEME_TEXT_COLOR
            p2.font.color.rgb = THEME_TITLE_COLOR
            p2.alignment = PP_ALIGN.RIGHT

        _add_shadow(lyrics_box)

    _add_hymn_background(slide)
    return slide


def _add_scripture_slides(prs, reference: str, translation: str):
    """Add scripture reading slides to the presentation."""
    if not reference or not reference.strip():
        return

    data = fetch_scripture(reference, translation)
    if not data or not data.get('slides'):
        return

    slides = data['slides']
    trans_name = data.get('translation_name', translation)

    for i, slide_data in enumerate(slides):
        is_last = (i == len(slides) - 1)
        if i == 0:
            _create_scripture_first_slide(prs, reference, trans_name, slide_data)
        else:
            _create_scripture_continuation_slide(prs, slide_data, is_last, trans_name)


def _add_theme_slide(prs, theme_image_path: Path | None):
    """Add a theme/separator slide using the uploaded theme image."""
    if theme_image_path and theme_image_path.exists():
        _add_full_image_slide(prs, theme_image_path)
    else:
        # Blank slide as fallback
        prs.slides.add_slide(prs.slide_layouts[6])


def generate_slides(order: OrderOfWorship, theme_name: str = None) -> Path:
    """Generate a PowerPoint presentation and return the file path."""
    _load_theme(theme_name)
    prs = Presentation()

    # Set slide dimensions
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Theme image path
    theme_path = None
    if order.themeImageFilename:
        candidate = OUTPUT_DIR / order.themeImageFilename
        if candidate.exists():
            theme_path = candidate

    # ===== SLIDE ORDER (from VBA Main()) =====

    # 1. Theme image (full-bleed)
    _add_theme_slide(prs, theme_path)

    # 2. Announcements slide (static image)
    _add_full_image_slide(prs, ANNOUNCEMENTS_SLIDE)

    # 3. Concerns slide (static image)
    _add_full_image_slide(prs, CONCERNS_SLIDE)

    # 4. Theme slide (separator)
    _add_theme_slide(prs, theme_path)

    # 5. Praise Hymn 1 (Opening Hymn)
    if order.praiseHymn1:
        _add_hymn_slides(prs, order.praiseHymn1, bg_type='hymn')

    # 6. Offering slide (static image)
    _add_full_image_slide(prs, OFFERING_SLIDE)

    # 7. Praise Hymn 2 (Offertory Hymn)
    if order.praiseHymn2:
        _add_hymn_slides(prs, order.praiseHymn2, bg_type='hymn')

    # 8. Doxology
    if order.doxology:
        _add_hymn_slides(prs, order.doxology, bg_type='hymn')

    # 9. Theme slide (separator)
    _add_theme_slide(prs, theme_path)

    # 10. Creed (uses creed background)
    if order.creed:
        _add_hymn_slides(prs, order.creed, bg_type='creed')

    # 11. Theme slide (separator)
    _add_theme_slide(prs, theme_path)

    # 12. Prayer Hymn
    if order.prayerHymn:
        _add_hymn_slides(prs, order.prayerHymn, bg_type='hymn')

    # 13. Prayer slide (static image)
    _add_full_image_slide(prs, PRAYER_SLIDE)

    # 14. Lord's Prayer / Liturgical Prayer (creed background)
    if order.liturgicalPrayer:
        _add_hymn_slides(prs, order.liturgicalPrayer, bg_type='creed')

    # 15. Theme slide (separator)
    _add_theme_slide(prs, theme_path)

    # 16. Scripture reading
    if order.scripture:
        _add_scripture_slides(prs, order.scripture, order.scriptureTranslation)
        _add_theme_slide(prs, theme_path)

    # 17. Closing Hymn
    if order.closingHymn:
        _add_hymn_slides(prs, order.closingHymn, bg_type='hymn')

    # 18. Theme slide (separator, final)
    _add_theme_slide(prs, theme_path)

    # Save
    OUTPUT_DIR.mkdir(exist_ok=True)
    filename = f"{order.date} - Slides.pptx"
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))
    return filepath
