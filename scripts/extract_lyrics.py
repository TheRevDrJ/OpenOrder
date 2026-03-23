"""
Extract lyrics and slide layout info from .pptx hymn files into JSON.

For each .pptx file, produces a JSON with:
- number: hymn number (parsed from filename)
- title: hymn title (parsed from filename)
- source: which collection it came from
- slides: list of slides, each with lines of text and line count

The extraction is smart about separating metadata (title, number, attribution,
offscreen placeholders) from actual lyrics/liturgy text. It does this by:
1. Skipping shapes with negative top position (offscreen placeholders)
2. On first slides: identifying title, number box, and attribution shapes
   by their position and content, and only extracting lyrics from the
   remaining text shapes
3. On continuation slides: extracting all visible text (which is just lyrics)

Skips Dropbox "conflicted copy" files.
"""

import json
import os
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


# Folders to process, with their source labels
FOLDERS = {
    "Hymnal/Hymns by Number": "umh",
    "Hymnal/Other General Services and Acts of Worship": "umh-services",
    "Hymnal/General Services": "umh-general",
    "Hymnal/Psalter": "umh-psalter",
    "The Faith We Sing/Number Index (Cross & Flame background)": "tfws",
}


def parse_hymn_filename(filename: str, source: str) -> dict:
    """Extract hymn number and title from filename."""
    name = Path(filename).stem

    if source == "umh":
        # Format: "057 - O For a Thousand Tongues to Sing"
        match = re.match(r"^(\d+)\s*-\s*(.+)$", name)
        if match:
            return {"number": match.group(1), "title": match.group(2).strip()}
    elif source in ("umh-services", "umh-general"):
        # Format: "881-The Apsotles' Creed, Traditional Version"
        match = re.match(r"^(\d+)-(.+)$", name)
        if match:
            return {"number": match.group(1), "title": match.group(2).strip()}
    elif source == "umh-psalter":
        # Format: "Psalm 001" or "Psalm 009 (11-20)"
        match = re.match(r"^Psalm\s+(\d+)(?:\s*\((.+)\))?$", name)
        if match:
            num = match.group(1)
            suffix = f" ({match.group(2)})" if match.group(2) else ""
            return {"number": num, "title": f"Psalm {num}{suffix}"}
    elif source == "tfws":
        # Format: "2001 - Grace Alone" or similar
        match = re.match(r"^(\d+)\s*-\s*(.+)$", name)
        if match:
            return {"number": match.group(1), "title": match.group(2).strip()}
        # Some might just be number-title without space
        match = re.match(r"^(\d+)(.+)$", name)
        if match:
            return {"number": match.group(1), "title": match.group(2).strip().lstrip("- ")}

    # Fallback
    return {"number": "0", "title": name}


def _is_metadata_text(text: str, hymn_number: str, hymn_title: str,
                       is_first_slide: bool = True) -> bool:
    """Check if a line of text is metadata rather than lyrics.

    NOTE: We no longer filter by title match here. Title/number/attribution
    are handled at the SHAPE level (by position), not the text level.
    This prevents stripping legitimate lyric lines that match the title
    (e.g., "Give Thanks" in TFWS 2036).

    This function now only filters out content that can appear WITHIN
    the lyrics shape itself (copyright lines, verse labels, etc.)
    """
    stripped = text.strip()
    if not stripped:
        return True

    # Copyright lines (these DO appear inside lyrics shapes)
    if stripped.startswith('©') or stripped.startswith('\u00a9'):
        return True

    # "FROM THE RITUAL OF..." type attribution (appears in services lyrics shapes)
    if is_first_slide and stripped.startswith("FROM THE"):
        return True

    # "FORMER METHODIST CHURCH" and similar (appears in services lyrics shapes)
    if is_first_slide and stripped.startswith("FORMER ") and stripped.isupper():
        return True

    return False


def _shape_is_offscreen(shape) -> bool:
    """Check if a shape is positioned offscreen (negative top)."""
    return shape.top is not None and shape.top < 0


def _shape_is_number_box(shape, hymn_number: str) -> bool:
    """Check if a shape is the hymn number box (far right, contains just the number)."""
    if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    # Number box is typically at left >= 9 inches and contains just the number
    if shape.left is not None and shape.left >= Inches(8):
        if text == hymn_number or text == hymn_number.lstrip("0"):
            return True
    return False


def _shape_is_title(shape) -> bool:
    """Check if a shape is the title text box (top ~0.33 inches, left ~1.67 inches)."""
    if shape.top is not None and shape.left is not None:
        # Title is at top of slide (< 1 inch from top) and left-aligned (< 3 inches from left)
        if shape.top < Inches(1) and shape.left < Inches(3):
            return True
    return False


def _shape_is_attribution(shape) -> bool:
    """Check if a shape contains attribution text (WORDS:, MUSIC:, etc.)."""
    if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    return text.startswith("WORDS:") or text.startswith("MUSIC:")


def extract_slide_text_smart(slide, slide_index: int, hymn_number: str,
                              hymn_title: str) -> list[str]:
    """
    Extract lyrics text from a slide using POSITION-BASED shape identification.

    The source .pptx slides have a consistent layout:
    - Offscreen placeholder: top < 0 (negative) -- SKIP
    - Title/attribution shape: top < 1.5 inches, left < 4 inches -- SKIP on slide 1
    - Number box: top < 1 inch, left > 8 inches -- SKIP on slide 1
    - Lyrics shape: top >= 1.5 inches, largest height -- KEEP

    On first slides (index 0), we ONLY extract from the lyrics shape,
    identified by position (top >= 1.5") and having the largest height.
    This avoids any text-matching heuristics that could accidentally
    strip legitimate lyrics (e.g., "Give Thanks" matching the title).

    On continuation slides (index > 0), we take all visible text since
    they contain only lyrics.
    """
    lines = []

    # Collect visible text shapes with their positions
    text_shapes = []
    for shape in slide.shapes:
        if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
            continue
        if _shape_is_offscreen(shape):
            continue
        text_shapes.append(shape)

    if slide_index == 0:
        # First slide: find the lyrics shape by position.
        # It's the shape with top >= 1.5 inches and the largest height.
        # This cleanly separates title/number/attribution from lyrics
        # without any text-matching that could strip valid lyric lines.
        lyrics_candidates = []
        for shape in text_shapes:
            top_in = (shape.top or 0) / 914400
            height_in = (shape.height or 0) / 914400
            if top_in >= 1.5:
                lyrics_candidates.append((shape, height_in))

        if lyrics_candidates:
            # Sort by height descending — lyrics shape is the tallest
            lyrics_candidates.sort(key=lambda x: x[1], reverse=True)
            lyrics_shape = lyrics_candidates[0][0]
            _extract_lines_from_shape(lyrics_shape, lines, hymn_number, hymn_title,
                                       is_first_slide=True)
        else:
            # Fallback: no shape found below 1.5". Take the tallest visible shape.
            tallest = max(text_shapes, key=lambda s: (s.height or 0), default=None)
            if tallest:
                _extract_lines_from_shape(tallest, lines, hymn_number, hymn_title,
                                           is_first_slide=True)
    else:
        # Continuation slides: take all visible text (it's just lyrics)
        for shape in text_shapes:
            _extract_lines_from_shape(shape, lines, hymn_number, hymn_title,
                                       is_first_slide=False)

    return lines


def _extract_lines_from_shape(shape, lines: list, hymn_number: str,
                               hymn_title: str, is_first_slide: bool = True):
    """Extract text lines from a shape, filtering metadata lines."""
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        # Check for soft returns (vertical tabs) — split them into separate lines
        if '\x0b' in text:
            sub_lines = text.split('\x0b')
            for sub in sub_lines:
                sub = sub.strip()
                if sub and not _is_metadata_text(sub, hymn_number, hymn_title,
                                                  is_first_slide):
                    lines.append(sub)
        else:
            # Keep "Refrain" as a label
            if text.lower() == "refrain":
                lines.append(text)
            elif not _is_metadata_text(text, hymn_number, hymn_title,
                                        is_first_slide):
                lines.append(text)


def process_pptx(filepath: str, source: str) -> dict | None:
    """Extract lyrics data from a single .pptx file."""
    filename = os.path.basename(filepath)

    # Skip conflicted copies
    if "conflicted copy" in filename:
        return None

    info = parse_hymn_filename(filename, source)

    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"  ERROR reading {filename}: {e}", file=sys.stderr)
        return None

    slides = []
    for slide_index, slide in enumerate(prs.slides):
        lines = extract_slide_text_smart(
            slide, slide_index, info["number"], info["title"]
        )
        if lines:
            slides.append({
                "lines": lines,
                "line_count": len(lines),
            })

    if not slides:
        return None

    return {
        "number": info["number"],
        "title": info["title"],
        "source": source,
        "filename": filename,
        "slide_count": len(slides),
        "slides": slides,
    }


def main():
    script_dir = Path(__file__).parent
    # Accept optional command-line argument for root directory
    if len(sys.argv) > 1:
        root_dir = Path(sys.argv[1])
    else:
        root_dir = script_dir.parent / "Hymnals"
    output_dir = script_dir.parent / "hymnal-json"
    output_dir.mkdir(exist_ok=True)

    # Create subdirs for each source
    for source in set(FOLDERS.values()):
        (output_dir / source).mkdir(exist_ok=True)

    total = 0
    converted = 0
    skipped = 0
    failed = 0

    for folder_rel, source in FOLDERS.items():
        folder = root_dir / folder_rel
        if not folder.exists():
            print(f"WARNING: Folder not found: {folder}")
            continue

        pptx_files = sorted(folder.glob("*.pptx"))
        print(f"\n{source}: {len(pptx_files)} .pptx files in {folder_rel}")

        for filepath in pptx_files:
            total += 1

            if "conflicted copy" in filepath.name:
                skipped += 1
                continue

            result = process_pptx(str(filepath), source)
            if result is None:
                failed += 1
                continue

            # Save JSON
            json_name = filepath.stem + ".json"
            json_path = output_dir / source / json_name
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(result, f, indent=2, ensure_ascii=False)

            converted += 1
            if converted % 50 == 0:
                print(f"  ... {converted} extracted so far")

    # Also build a master index
    index = []
    for source_dir in output_dir.iterdir():
        if source_dir.is_dir():
            for json_file in sorted(source_dir.glob("*.json")):
                with open(json_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                index.append({
                    "number": data["number"],
                    "title": data["title"],
                    "source": data["source"],
                    "slide_count": data["slide_count"],
                    "file": f"{source_dir.name}/{json_file.name}",
                })

    index.sort(key=lambda x: (x["source"], x["number"].zfill(6)))
    index_path = output_dir / "index.json"
    with open(index_path, "w", encoding="utf-8") as f:
        json.dump(index, f, indent=2, ensure_ascii=False)

    print(f"\n=== Done ===")
    print(f"  Converted: {converted}")
    print(f"  Skipped: {skipped}")
    print(f"  Failed: {failed}")
    print(f"  Total: {total}")
    print(f"  Index: {index_path} ({len(index)} entries)")


if __name__ == "__main__":
    main()
